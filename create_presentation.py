#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tao PowerPoint Presentation cho Biometric Security Lab
Su dung thu vien python-pptx
"""

import sys
import os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Cấu hình màu
COLOR_PRIMARY = RGBColor(102, 126, 234)      # Xanh
COLOR_SECONDARY = RGBColor(118, 75, 162)    # Tím
COLOR_DANGER = RGBColor(220, 53, 69)        # Đỏ
COLOR_SUCCESS = RGBColor(40, 167, 69)        # Xanh lá
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK = RGBColor(51, 51, 51)
COLOR_LIGHT_BG = RGBColor(248, 249, 250)

def add_title_slide(prs, title, subtitle="", author="", date=""):
    """Tạo slide tiêu đề"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    # Author and Date
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True

    if author:
        p = info_frame.add_paragraph()
        p.text = f"Sinh viên: {author}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)

    if date:
        p = info_frame.add_paragraph()
        p.text = f"Ngày: {date}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items):
    """Tạo slide nội dung với tiêu đề và bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # Title underline
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(2), Inches(0))
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(3)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item["text"]
        p.level = item.get("level", 0)
        p.font.size = Pt(item.get("size", 18))
        p.font.color.rgb = item.get("color", COLOR_DARK)
        p.font.bold = item.get("bold", False)
        p.space_before = Pt(item.get("space_before", 6))
        p.space_after = Pt(item.get("space_after", 6))

    return slide


def add_two_column_slide(prs, title, left_items, right_items):
    """Tạo slide với 2 cột"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # Title underline
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(2), Inches(0))
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(3)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = item["text"]
        p.level = item.get("level", 0)
        p.font.size = Pt(item.get("size", 16))
        p.font.color.rgb = item.get("color", COLOR_DARK)
        p.font.bold = item.get("bold", False)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = item["text"]
        p.level = item.get("level", 0)
        p.font.size = Pt(item.get("size", 16))
        p.font.color.rgb = item.get("color", COLOR_DARK)
        p.font.bold = item.get("bold", False)

    return slide


def create_presentation():
    """Tạo toàn bộ presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("[*] Đang tạo PowerPoint...")

    # SLIDE 1: TITLE
    print("  [+] Slide 1: Tiêu đề")
    add_title_slide(prs,
                   "🔐 Biometric Security Lab",
                   "Tấn Công & Phòng Vệ Sinh Trắc Học Khuôn Mặt",
                   "[Tên của bạn]",
                   "20/04/2026")

    # SLIDE 2: OVERVIEW
    print("  [+] Slide 2: Tổng quan")
    add_content_slide(prs, "📋 Tổng Quan", [
        {"text": "Vấn Đề", "size": 26, "bold": True, "color": COLOR_SECONDARY},
        {"text": "Sinh trắc học khuôn mặt dễ bị tấn công nếu không có bảo vệ đúng", "level": 1, "size": 16},
        {"text": "Nhiều ứng dụng thực tế không kiểm tra Liveness Detection", "level": 1, "size": 16},
        {"text": "Attacker có thể dùng ảnh tĩnh để bypass xác thực", "level": 1, "size": 16},
        {"text": "Giải Pháp Của Em", "size": 26, "bold": True, "color": COLOR_SECONDARY, "space_before": 20},
        {"text": "✅ Xây dựng server xác thực khuôn mặt đơn giản", "level": 1, "size": 16},
        {"text": "✅ Chứng minh 2 loại tấn công: Reconnaissance + Injection", "level": 1, "size": 16, "color": COLOR_DANGER},
        {"text": "✅ Đề xuất giải pháp bảo mật", "level": 1, "size": 16},
    ])

    # SLIDE 3: ARCHITECTURE
    print("  [+] Slide 3: Kiến trúc")
    add_content_slide(prs, "🏗️ Kiến Trúc Hệ Thống", [
        {"text": "Flask Server (app.py)", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": "📤 POST /register → Đăng ký khuôn mặt", "level": 1, "size": 15},
        {"text": "🔍 POST /authenticate → Xác thực khuôn mặt", "level": 1, "size": 15},
        {"text": "📋 GET /users → Liệt kê user (⚠️ NO AUTH)", "level": 1, "size": 15, "color": COLOR_DANGER},
        {"text": "Face Recognition Library", "size": 20, "bold": True, "color": COLOR_PRIMARY, "space_before": 16},
        {"text": "Trích xuất vector 128-D từ khuôn mặt", "level": 1, "size": 15},
        {"text": "So khớp bằng Euclidean distance", "level": 1, "size": 15},
        {"text": "Database (pickle file)", "size": 20, "bold": True, "color": COLOR_PRIMARY, "space_before": 16},
        {"text": "Lưu: User name + face encoding + registered time", "level": 1, "size": 15},
    ])

    # SLIDE 4: RECONNAISSANCE
    print("  [+] Slide 4: Trinh sát")
    add_content_slide(prs, "🕵️ Tấn Công 1: Trinh Sát", [
        {"text": "Cách Hoạt Động", "size": 26, "bold": True, "color": COLOR_SECONDARY},
        {"text": "Attacker gửi GET /users (không cần auth)", "level": 1, "size": 16},
        {"text": "Server trả về JSON chứa danh sách tất cả user", "level": 1, "size": 16},
        {"text": "Attacker lấy được tên nạn nhân", "level": 1, "size": 16},
        {"text": "Vấn Đề Bảo Mật", "size": 26, "bold": True, "color": COLOR_SECONDARY, "space_before": 20},
        {"text": "❌ Information Disclosure", "level": 0, "size": 18, "bold": True, "color": COLOR_DANGER},
        {"text": "Endpoint không yêu cầu authentication", "level": 1, "size": 16},
        {"text": "Bất kỳ ai cũng có thể gọi GET /users", "level": 1, "size": 16},
    ])

    # SLIDE 5: INJECTION ATTACK
    print("  [+] Slide 5: Injection attack")
    add_content_slide(prs, "💉 Tấn Công 2: Injection Attack", [
        {"text": "Quy Trình Tấn Công", "size": 26, "bold": True, "color": COLOR_SECONDARY},
        {"text": "Lấy ảnh khuôn mặt nạn nhân (từ mạng)", "level": 1, "size": 16},
        {"text": "Encode ảnh thành base64", "level": 1, "size": 16},
        {"text": "POST trực tiếp tới /authenticate", "level": 1, "size": 16},
        {"text": "Server xác thực thành công ❌", "level": 1, "size": 16, "color": COLOR_DANGER},
        {"text": "Tại Sao Thành Công?", "size": 26, "bold": True, "color": COLOR_SECONDARY, "space_before": 20},
        {"text": "❌ Không có Liveness Detection", "level": 0, "size": 18, "bold": True, "color": COLOR_DANGER},
        {"text": "Server không kiểm tra ảnh có phải từ khuôn mặt sống thật", "level": 1, "size": 15},
    ])

    # SLIDE 6: DEMO RESULTS
    print("  [+] Slide 6: Kết quả demo")
    add_content_slide(prs, "🎯 Kết Quả Demo", [
        {"text": "Trinh Sát", "size": 24, "bold": True, "color": COLOR_PRIMARY},
        {"text": "[+] Tìm thấy 1 user: Nguyen Van A", "level": 0, "size": 14, "bold": False},
        {"text": "Injection Attack - 1 Ảnh", "size": 24, "bold": True, "color": COLOR_PRIMARY, "space_before": 16},
        {"text": "[!!!] INJECTION THÀNH CÔNG!", "level": 0, "size": 16, "bold": True, "color": COLOR_SUCCESS},
        {"text": "Xác thực là: Nguyen Van A | Confidence: 0.92", "level": 0, "size": 14},
        {"text": "Brute-Force - 10 Ảnh", "size": 24, "bold": True, "color": COLOR_PRIMARY, "space_before": 16},
        {"text": "Tổng thử: 10 | Thành công: 6 | Thất bại: 4", "level": 0, "size": 14, "bold": True, "color": COLOR_DANGER},
        {"text": "Tỷ lệ bypass: 60%", "level": 0, "size": 16, "bold": True, "color": COLOR_DANGER},
    ])

    # SLIDE 7: RECOMMENDATIONS
    print("  [+] Slide 7: Khuyến nghị")
    add_content_slide(prs, "🛡️ Khuyến Nghị Bảo Mật", [
        {"text": "1. Liveness Detection", "size": 22, "bold": True, "color": COLOR_SECONDARY},
        {"text": "Yêu cầu chuyển động (nháy mắt, cười, quay đầu)", "level": 1, "size": 15},
        {"text": "Phát hiện ảnh 2D vs khuôn mặt sống (dùng AI/ML)", "level": 1, "size": 15},
        {"text": "2. Authentication & Authorization", "size": 22, "bold": True, "color": COLOR_SECONDARY, "space_before": 14},
        {"text": "Endpoint /users cần JWT token (chỉ admin)", "level": 1, "size": 15},
        {"text": "3. Rate Limiting & Logging", "size": 22, "bold": True, "color": COLOR_SECONDARY, "space_before": 14},
        {"text": "Giới hạn request (5 lần/phút)", "level": 1, "size": 15},
        {"text": "Ghi log tất cả xác thực + phát hiện pattern tấn công", "level": 1, "size": 15},
    ])

    # SLIDE 8: CONCLUSION
    print("  [+] Slide 8: Kết luận")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "🎓 Kết Luận"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    conclusions = [
        "✅ Sinh trắc học khuôn mặt là tốt nhưng cần bảo vệ bổ sung",
        "✅ Liveness Detection là bước quan trọng để chống tấn công",
        "✅ Authentication trên tất cả endpoint là cần thiết",
        "✅ Bảo mật sinh trắc học = kết hợp nhiều lớp bảo vệ",
    ]

    for i, conclusion in enumerate(conclusions):
        p = content_frame.paragraphs[0] if i == 0 else content_frame.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(8)

    # Thank you
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    thank_frame = thank_box.text_frame
    p = thank_frame.paragraphs[0]
    p.text = "Cảm ơn! 🙏"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = "C:\\Users\\ADMIN\\biometric-security-lab\\Biometric_Security_Lab_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ PowerPoint tạo thành công!")
    print(f"📁 Đường dẫn: {output_path}")
    print(f"📊 Tổng slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
