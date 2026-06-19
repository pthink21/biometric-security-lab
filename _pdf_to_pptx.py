"""Convert the light-mode PDF into a .pptx where each page becomes a full-bleed image slide."""
import fitz
import io
from pptx import Presentation
from pptx.util import Emu

SRC_PDF = r"D:\Biometrics_Presentation_LightMode.pdf"
DST_PPTX = r"D:\Biometrics_Presentation_LightMode.pptx"
DPI = 200

def main():
    doc = fitz.open(SRC_PDF)
    prs = Presentation()

    page0 = doc[0]
    pt_w, pt_h = page0.rect.width, page0.rect.height
    prs.slide_width = Emu(int(pt_w / 72.0 * 914400))
    prs.slide_height = Emu(int(pt_h / 72.0 * 914400))

    blank = prs.slide_layouts[6]
    zoom = DPI / 72.0
    mat = fitz.Matrix(zoom, zoom)

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        buf = io.BytesIO(pix.tobytes("jpeg", jpg_quality=90))

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)
        print(f"  - Slide {i+1}/{len(doc)}")

    prs.save(DST_PPTX)
    doc.close()
    print(f"[+] Saved: {DST_PPTX}")

if __name__ == "__main__":
    main()
